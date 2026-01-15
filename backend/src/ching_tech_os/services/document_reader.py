"""
Document Reader Service

提供文件內容讀取功能，支援 Word、Excel、PowerPoint、PDF 等文件格式的文字提取。
"""

import os
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


# 支援的副檔名
SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".pdf"}

# 舊版格式（提示轉檔）
LEGACY_EXTENSIONS = {".doc", ".xls", ".ppt"}

# 檔案大小限制（bytes）
MAX_FILE_SIZE = {
    ".pdf": 10 * 1024 * 1024,   # 10 MB
    ".docx": 10 * 1024 * 1024,  # 10 MB
    ".xlsx": 5 * 1024 * 1024,   # 5 MB
    ".pptx": 10 * 1024 * 1024,  # 10 MB
}

# 最大輸出字元數
MAX_OUTPUT_CHARS = 100000


class DocumentReadError(Exception):
    """文件讀取錯誤"""
    pass


class PasswordProtectedError(DocumentReadError):
    """文件有密碼保護"""
    pass


class CorruptedFileError(DocumentReadError):
    """文件損壞"""
    pass


class UnsupportedFormatError(DocumentReadError):
    """不支援的格式"""
    pass


class FileTooLargeError(DocumentReadError):
    """檔案過大"""
    pass


@dataclass
class DocumentContent:
    """文件解析結果"""
    text: str                      # 提取的純文字內容
    format: str                    # 原始格式 (docx, xlsx, pptx, pdf)
    page_count: Optional[int]      # 頁數（PDF）或工作表數（Excel）
    metadata: dict                 # 額外資訊（標題、作者等）
    truncated: bool                # 是否因大小限制而截斷
    error: Optional[str]           # 解析錯誤訊息（部分成功時）


def get_supported_extensions() -> set[str]:
    """取得支援的副檔名"""
    return SUPPORTED_EXTENSIONS.copy()


def is_supported(file_path: str) -> bool:
    """檢查檔案是否支援"""
    ext = Path(file_path).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


def is_legacy_format(file_path: str) -> bool:
    """檢查是否為舊版格式"""
    ext = Path(file_path).suffix.lower()
    return ext in LEGACY_EXTENSIONS


def extract_text(file_path: str) -> DocumentContent:
    """
    提取文件的文字內容

    Args:
        file_path: 檔案路徑

    Returns:
        DocumentContent 包含提取的文字和元資料

    Raises:
        FileNotFoundError: 檔案不存在
        UnsupportedFormatError: 不支援的格式
        FileTooLargeError: 檔案過大
        PasswordProtectedError: 文件有密碼保護
        CorruptedFileError: 文件損壞
    """
    path = Path(file_path)

    # 檢查檔案存在
    if not path.exists():
        raise FileNotFoundError(f"檔案不存在: {file_path}")

    ext = path.suffix.lower()

    # 檢查舊版格式
    if ext in LEGACY_EXTENSIONS:
        raise UnsupportedFormatError(
            f"不支援舊版格式 {ext}，請轉存為新版格式 "
            f"(.docx/.xlsx/.pptx)"
        )

    # 檢查支援的格式
    if ext not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFormatError(
            f"不支援的檔案格式: {ext}。"
            f"支援的格式: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    # 檢查檔案大小
    file_size = path.stat().st_size
    max_size = MAX_FILE_SIZE.get(ext, 10 * 1024 * 1024)
    if file_size > max_size:
        raise FileTooLargeError(
            f"檔案過大 ({file_size / 1024 / 1024:.1f} MB)，"
            f"請使用小於 {max_size / 1024 / 1024:.0f} MB 的檔案"
        )

    # 根據格式選擇解析器
    extractors = {
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".pdf": _extract_pdf,
    }

    extractor = extractors[ext]
    return extractor(file_path)


def _extract_docx(file_path: str) -> DocumentContent:
    """解析 Word 文件"""
    try:
        doc = Document(file_path)
    except Exception as e:
        error_msg = str(e).lower()
        if "password" in error_msg or "encrypted" in error_msg:
            raise PasswordProtectedError("此文件有密碼保護，無法讀取")
        raise CorruptedFileError(f"無法解析 Word 文件: {e}")

    paragraphs = []

    # 提取段落
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # 提取表格
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                paragraphs.append(row_text)

    text = "\n".join(paragraphs)

    # 提取元資料
    metadata = {}
    core_props = doc.core_properties
    if core_props.title:
        metadata["title"] = core_props.title
    if core_props.author:
        metadata["author"] = core_props.author

    # 檢查是否需要截斷
    truncated = False
    if len(text) > MAX_OUTPUT_CHARS:
        text = text[:MAX_OUTPUT_CHARS] + f"\n\n[內容已截斷，原文共 {len(text)} 字元]"
        truncated = True

    return DocumentContent(
        text=text,
        format="docx",
        page_count=None,
        metadata=metadata,
        truncated=truncated,
        error=None
    )


def _extract_xlsx(file_path: str) -> DocumentContent:
    """解析 Excel 試算表"""
    try:
        wb = load_workbook(file_path, data_only=True, read_only=True)
    except Exception as e:
        error_msg = str(e).lower()
        if "password" in error_msg or "encrypted" in error_msg:
            raise PasswordProtectedError("此文件有密碼保護，無法讀取")
        raise CorruptedFileError(f"無法解析 Excel 檔案: {e}")

    result = []
    sheet_count = len(wb.sheetnames)

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        result.append(f"=== 工作表: {sheet_name} ===")

        row_count = 0
        for row in sheet.iter_rows(values_only=True):
            # 過濾完全空白的行
            if any(cell is not None for cell in row):
                row_text = " | ".join(
                    str(cell) if cell is not None else ""
                    for cell in row
                )
                result.append(row_text)
                row_count += 1

        if row_count == 0:
            result.append("（空白工作表）")

        result.append("")  # 工作表間空行

    wb.close()

    text = "\n".join(result)

    # 檢查是否需要截斷
    truncated = False
    if len(text) > MAX_OUTPUT_CHARS:
        text = text[:MAX_OUTPUT_CHARS] + f"\n\n[內容已截斷，原文共 {len(text)} 字元]"
        truncated = True

    return DocumentContent(
        text=text,
        format="xlsx",
        page_count=sheet_count,
        metadata={"sheet_names": wb.sheetnames},
        truncated=truncated,
        error=None
    )


def _extract_pptx(file_path: str) -> DocumentContent:
    """解析 PowerPoint 簡報"""
    try:
        prs = Presentation(file_path)
    except Exception as e:
        error_msg = str(e).lower()
        if "password" in error_msg or "encrypted" in error_msg:
            raise PasswordProtectedError("此文件有密碼保護，無法讀取")
        raise CorruptedFileError(f"無法解析 PowerPoint 檔案: {e}")

    result = []
    slide_count = len(prs.slides)

    for i, slide in enumerate(prs.slides, 1):
        result.append(f"=== 投影片 {i} ===")

        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

        if slide_text:
            result.extend(slide_text)
        else:
            result.append("（無文字內容）")

        result.append("")  # 投影片間空行

    text = "\n".join(result)

    # 檢查是否需要截斷
    truncated = False
    if len(text) > MAX_OUTPUT_CHARS:
        text = text[:MAX_OUTPUT_CHARS] + f"\n\n[內容已截斷，原文共 {len(text)} 字元]"
        truncated = True

    return DocumentContent(
        text=text,
        format="pptx",
        page_count=slide_count,
        metadata={},
        truncated=truncated,
        error=None
    )


def _extract_pdf(file_path: str) -> DocumentContent:
    """解析 PDF 文件"""
    try:
        with fitz.open(file_path) as doc:
            # 檢查是否需要密碼
            if doc.needs_pass:
                raise PasswordProtectedError("此文件有密碼保護，無法讀取")

            result = []
            page_count = len(doc)
            has_text = False

            for page in doc:
                text = page.get_text()
                if text.strip():
                    result.append(text)
                    has_text = True

            # 檢查是否為純圖片 PDF
            if not has_text:
                return DocumentContent(
                    text="此 PDF 為掃描圖片，沒有可提取的文字。建議截圖後上傳讓 AI 讀取。",
                    format="pdf",
                    page_count=page_count,
                    metadata={"is_scanned": True},
                    truncated=False,
                    error="純圖片 PDF，無文字層"
                )

            text = "\n".join(result)

            # 檢查是否需要截斷
            truncated = False
            if len(text) > MAX_OUTPUT_CHARS:
                text = text[:MAX_OUTPUT_CHARS] + f"\n\n[內容已截斷，原文共 {len(text)} 字元]"
                truncated = True

            return DocumentContent(
                text=text,
                format="pdf",
                page_count=page_count,
                metadata={},
                truncated=truncated,
                error=None
            )
    except PasswordProtectedError:
        raise
    except Exception as e:
        raise CorruptedFileError(f"無法解析 PDF 檔案: {e}")


@dataclass
class PdfConversionResult:
    """PDF 轉圖片結果"""
    success: bool
    total_pages: int               # PDF 總頁數
    converted_pages: int           # 實際轉換的頁數
    images: list[str]              # 轉換後的圖片路徑列表
    message: str                   # 人類可讀的結果描述
    error: Optional[str] = None    # 錯誤訊息（如有）


def _parse_pages_param(pages: str, total_pages: int) -> list[int]:
    """
    解析 pages 參數，回傳要轉換的頁面索引列表（0-based）

    Args:
        pages: 頁面參數，如 "0"、"1"、"1-3"、"1,3,5"、"all"
        total_pages: PDF 總頁數

    Returns:
        頁面索引列表（0-based）
    """
    if pages == "0":
        # 只查詢頁數，不轉換
        return []
    if pages == "all":
        return list(range(total_pages))

    result = []
    parts = pages.split(",")
    for part in parts:
        part = part.strip()
        if "-" in part:
            # 範圍格式：1-3
            start, end = part.split("-", 1)
            start_idx = int(start) - 1  # 轉成 0-based
            end_idx = int(end) - 1
            for i in range(start_idx, end_idx + 1):
                if 0 <= i < total_pages and i not in result:
                    result.append(i)
        else:
            # 單頁格式：1
            idx = int(part) - 1  # 轉成 0-based
            if 0 <= idx < total_pages and idx not in result:
                result.append(idx)

    return sorted(result)


def convert_pdf_to_images(
    file_path: str,
    output_dir: str,
    pages: str = "all",
    dpi: int = 150,
    output_format: str = "png",
    max_pages: int = 20
) -> PdfConversionResult:
    """
    將 PDF 轉換為圖片

    Args:
        file_path: PDF 檔案路徑
        output_dir: 輸出目錄路徑
        pages: 要轉換的頁面，"0" 只查詢頁數、"1"、"1-3"、"1,3,5"、"all"
        dpi: 解析度，預設 150
        output_format: 輸出格式，png 或 jpg
        max_pages: 最大頁數限制，預設 20

    Returns:
        PdfConversionResult 包含轉換結果

    Raises:
        FileNotFoundError: 檔案不存在
        UnsupportedFormatError: 不是 PDF 格式
        PasswordProtectedError: 文件有密碼保護
        CorruptedFileError: 文件損壞
    """
    path = Path(file_path)

    # 檢查檔案存在
    if not path.exists():
        raise FileNotFoundError(f"PDF 檔案不存在: {file_path}")

    # 檢查是否為 PDF
    if path.suffix.lower() != ".pdf":
        raise UnsupportedFormatError(f"檔案不是 PDF 格式: {path.suffix}")

    try:
        with fitz.open(file_path) as doc:
            # 檢查是否需要密碼
            if doc.needs_pass:
                raise PasswordProtectedError("此 PDF 有密碼保護，無法轉換")

            total_pages = len(doc)

            # 解析要轉換的頁面
            page_indices = _parse_pages_param(pages, total_pages)

            # 如果 pages="0"，只回傳頁數資訊
            if not page_indices:
                return PdfConversionResult(
                    success=True,
                    total_pages=total_pages,
                    converted_pages=0,
                    images=[],
                    message=f"此 PDF 共有 {total_pages} 頁"
                )

            # 限制最大頁數
            if len(page_indices) > max_pages:
                page_indices = page_indices[:max_pages]

            # 確保輸出目錄存在
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)

            # 轉換設定
            zoom = dpi / 72
            mat = fitz.Matrix(zoom, zoom)
            images = []

            for idx in page_indices:
                page = doc[idx]
                pix = page.get_pixmap(matrix=mat)

                # 輸出檔名（使用 1-based 頁碼）
                img_path = output_path / f"page-{idx + 1}.{output_format}"
                pix.save(str(img_path))
                images.append(str(img_path))

            # 組合結果訊息
            converted_count = len(images)
            if converted_count == total_pages:
                message = f"已將全部 {total_pages} 頁轉換為圖片"
            elif converted_count == 1:
                message = f"已將第 {page_indices[0] + 1} 頁轉換為圖片（共 {total_pages} 頁）"
            else:
                page_desc = ", ".join(str(i + 1) for i in page_indices)
                message = f"已將第 {page_desc} 頁轉換為圖片（共 {total_pages} 頁）"

            return PdfConversionResult(
                success=True,
                total_pages=total_pages,
                converted_pages=converted_count,
                images=images,
                message=message
            )

    except PasswordProtectedError:
        raise
    except UnsupportedFormatError:
        raise
    except FileNotFoundError:
        raise
    except Exception as e:
        raise CorruptedFileError(f"無法轉換 PDF 檔案: {e}")
